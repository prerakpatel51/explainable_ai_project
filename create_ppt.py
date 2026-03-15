#!/usr/bin/env python3
"""XAI Presentation — clean layout, no overlapping, proper spacing."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BASE = "/home1/ppatel2025/xai_project_2"
EVAL = f"{BASE}/output/evaluation"
XAI_OUT = f"{BASE}/output/xai"
DOCS = f"{BASE}/docs/images"

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
CARD_BG = RGBColor(0x25, 0x25, 0x45)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xBB, 0xBB, 0xBB)
ORANGE = RGBColor(0xFF, 0x9F, 0x1C)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
PURPLE = RGBColor(0xA2, 0x9B, 0xFE)

def add_bg(s):
    f = s.background.fill; f.solid(); f.fore_color.rgb = DARK_BG

def tb(s, l, t, w, h, txt, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a
    return tf

def title(s, txt):
    tb(s, 0.8, 0.3, 11.7, 0.7, txt, 36, ACCENT, True)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.7), Pt(2))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

def card(s, l, t, w, h, bc=ACCENT):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD_BG; sh.line.color.rgb = bc; sh.line.width = Pt(1.5)
    return sh

def bullets(s, l, t, w, h, items, sz=18):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        txt, lvl, clr = item[0], item[1], item[2]
        bld = item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz - lvl * 2); p.font.color.rgb = clr or WHITE; p.font.bold = bld; p.level = lvl; p.space_after = Pt(4)
    return tf

def img(s, path, l, t, w=None, h=None):
    if os.path.exists(path):
        kw = {}
        if w: kw['width'] = Inches(w)
        if h: kw['height'] = Inches(h)
        s.shapes.add_picture(path, Inches(l), Inches(t), **kw)
        return True
    return False

def snum(s, n):
    tb(s, 12.3, 7.05, 0.8, 0.35, str(n), 11, LGRAY, False, PP_ALIGN.RIGHT)

# Filenames mapping for each class across models/domains
# real_model/real, real_model/sketch, sketch_model/real, sketch_model/sketch
CLASS_FILES = {
    'ladder':     ('ladder_0560.png',     'ladder_0544.png',     'ladder_0560.png',     'ladder_0544.png'),
    'lion':       ('lion_0658.png',       'lion_0643.png',       'lion_0658.png',       'lion_0643.png'),
    'tiger':      ('tiger_1405.png',      'tiger_1381.png',      'tiger_1405.png',      'tiger_1381.png'),
    'bicycle':    ('bicycle_0201.png',    'bicycle_0195.png',    'bicycle_0201.png',    'bicycle_0195.png'),
    'lighthouse': ('lighthouse_0638.png', 'lighthouse_0623.png', 'lighthouse_0638.png', 'lighthouse_0623.png'),
    'banana':     ('banana_0100.png',     'banana_0094.png',     'banana_0100.png',     'banana_0094.png'),
    'spider':     ('spider_1165.png',     'spider_1141.png',     'spider_1165.png',     'spider_1141.png'),
    'scissors':   ('scissors_0865.png',   'scissors_0846.png',   'scissors_0865.png',   'scissors_0846.png'),
    'saxophone':  ('saxophone_0831.png',  'saxophone_0801.png',  'saxophone_0831.png',  'saxophone_0801.png'),
    'snake':      ('snake_1085.png',     'snake_1061.png',      'snake_1085.png',      'snake_1061.png'),
    'umbrella':   ('umbrella_1505.png',  'umbrella_1480.png',   'umbrella_1505.png',    'umbrella_1480.png'),
    'airplane':   ('airplane_0040.png',  'airplane_0034.png',   'airplane_0040.png',    'airplane_0034.png'),
}

METHODS = [("Grad-CAM", "gradcam"), ("Grad-CAM++", "gradcam_pp"), ("Integrated Gradients", "integrated_gradients"), ("LIME", "lime")]

def get_img_path(model, method_dir, domain, cls_name):
    """Get image path: model=real_model|sketch_model, domain=real|sketch"""
    if model == 'real_model' and domain == 'real': fi = 0
    elif model == 'real_model' and domain == 'sketch': fi = 1
    elif model == 'sketch_model' and domain == 'real': fi = 2
    else: fi = 3
    fname = CLASS_FILES.get(cls_name, (None,None,None,None))[fi]
    if fname:
        return f"{XAI_OUT}/{model}/attribution_maps/{method_dir}/{domain}/{fname}"
    return ""

def add_2x2_grid(s, model, domain, cls_name, y_start=1.7):
    """Add a 2x2 grid of attribution maps with constrained sizes — no overlapping."""
    iw = 4.8; ih = 1.75; lbl_h = 0.28; gap_x = 0.3; gap_y = 0.1
    total_w = 2 * iw + gap_x
    start_x = (13.333 - total_w) / 2
    positions = [(0, 0), (0, 1), (1, 0), (1, 1)]
    for idx, (mname, mdir) in enumerate(METHODS):
        row, col = positions[idx]
        x = start_x + col * (iw + gap_x)
        y = y_start + row * (ih + lbl_h + gap_y)
        tb(s, x, y, iw, lbl_h, mname, 15, ACCENT, True, PP_ALIGN.CENTER)
        p = get_img_path(model, mdir, domain, cls_name)
        img(s, p, x, y + lbl_h, w=iw, h=ih)

def add_model_comparison(slide_counter, cls_name, domain, domain_label):
    """2 slides comparing Real Model vs Sketch Model on same image, 2 methods per slide."""
    pretty = cls_name.replace('_', ' ').title()
    for slide_idx, method_pair in enumerate([(0, 1), (2, 3)]):
        s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
        title(s, f"Model Comparison: {pretty} ({domain_label} Image)")
        pair_names = " & ".join([METHODS[i][0] for i in method_pair])
        tb(s, 0.5, 1.15, 12.3, 0.3, f"{pair_names} — How does each model explain the same image?", 16, LGRAY, False, PP_ALIGN.CENTER)

        iw = 5.2; ih = 2.0; gap_x = 0.4
        total_w = 2 * iw + gap_x
        start_x = (13.333 - total_w) / 2

        for col_idx, mi in enumerate(method_pair):
            mname, mdir = METHODS[mi]
            x = start_x + col_idx * (iw + gap_x)
            # Real Model row
            tb(s, x, 1.55, iw, 0.25, f"{mname} — Real Model", 14, RED, True, PP_ALIGN.CENTER)
            img(s, get_img_path('real_model', mdir, domain, cls_name), x, 1.82, w=iw, h=ih)
            # Sketch Model row
            tb(s, x, 3.95, iw, 0.25, f"{mname} — Sketch Model", 14, GREEN, True, PP_ALIGN.CENTER)
            img(s, get_img_path('sketch_model', mdir, domain, cls_name), x, 4.22, w=iw, h=ih)

        # Brief note
        note_y = 6.35
        if domain == 'sketch':
            note = "The sketch model (green) produces tighter, more shape-focused attributions. The real model (red) shows diffuse attention — searching for texture features absent in sketches."
        else:
            note = "Both models focus on similar regions on real images, but the sketch model's shape-based features produce cleaner, more consistent attributions even on photographs."
        c2 = card(s, 0.4, note_y, 12.5, 0.85, PURPLE)
        tf2 = c2.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0.15); tf2.margin_top = Inches(0.06)
        p2 = tf2.paragraphs[0]; p2.text = note; p2.font.size = Pt(14); p2.font.color.rgb = WHITE
        snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# SLIDE 1: Title
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.5, 1.0, 12.3, 1.2, "Explainability Under Domain Shift", 48, ACCENT, True, PP_ALIGN.CENTER)
tb(s, 0.5, 2.4, 12.3, 0.8, "Are Neural Network Explanations Consistent Across Visual Domains?", 26, WHITE, False, PP_ALIGN.CENTER)
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.5), Inches(7.3), Pt(2))
sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()
tb(s, 0.5, 3.8, 12.3, 0.5, "ResNet-152  |  DomainNet (Real & Sketch)  |  81 Classes  |  4 XAI Methods  |  4 Evaluation Axes", 18, LGRAY, False, PP_ALIGN.CENTER)
tb(s, 0.5, 5.0, 12.3, 0.5, "Prerak Patel", 28, ORANGE, True, PP_ALIGN.CENTER)
tb(s, 0.5, 5.6, 12.3, 0.5, "Florida Institute of Technology", 20, LGRAY, False, PP_ALIGN.CENTER)
snum(s, 1)

# ================================================================
# SLIDE 2: Motivation
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Motivation & Research Question")
bullets(s, 0.8, 1.3, 5.8, 3.5, [
    ("Why Explainability Matters", 0, ORANGE, True),
    ("Deep learning is a black box — high accuracy, no transparency", 1, WHITE),
    ("Safety-critical domains (medical, autonomous) need trust", 1, WHITE),
    ("Regulations (EU AI Act) increasingly require explanations", 1, WHITE),
    ("", 0, None),
    ("The Problem: Domain Shift", 0, ORANGE, True),
    ("Models trained on photos deployed on X-rays, sketches, satellite images", 1, WHITE),
    ("Even if accuracy holds, do explanations remain reliable?", 1, WHITE),
    ("A wrong explanation on a correct prediction is dangerous", 1, WHITE),
], 20)
bullets(s, 7.2, 1.3, 5.3, 3.5, [
    ("Experimental Design", 0, ACCENT, True),
    ("Train ResNet-152 on Real photographs", 1, WHITE),
    ("Train ResNet-152 on Sketch drawings", 1, WHITE),
    ("Test both models on BOTH domains", 1, WHITE),
    ("Apply 4 XAI methods to each combination", 1, WHITE),
    ("Evaluate explanations on 4 quantitative axes", 1, WHITE),
    ("Compare: which explanations survive domain shift?", 1, YELLOW),
], 20)
c = card(s, 0.8, 5.0, 11.7, 1.8, ACCENT)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.text = "Central Research Question"; p.font.size = Pt(18); p.font.color.rgb = ORANGE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = 'When a neural network is trained on one visual style, do its explanations remain consistent, faithful,\nand stable — even on images from a completely different visual domain?'; p.font.size = Pt(22); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
snum(s, 2)

# ================================================================
# SLIDE 3: Dataset
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Dataset: DomainNet Benchmark")
bullets(s, 0.8, 1.3, 5.8, 5.5, [
    ("Dataset Overview", 0, ORANGE, True),
    ("DomainNet: largest domain adaptation benchmark", 1, WHITE),
    ("We use 2 domains: Real (photos) & Sketch (drawings)", 1, WHITE),
    ("81 classes across 10 semantic categories", 1, WHITE),
    ("", 0, None),
    ("Data Splits (Stratified 80/10/10)", 0, ORANGE, True),
    ("Real:    29,501 train / 3,515 val / 3,515 test", 1, WHITE),
    ("Sketch: 29,501 train / 2,724 val / 2,724 test", 1, WHITE),
    ("", 0, None),
    ("Data Balancing", 0, ORANGE, True),
    ("Cross-domain oversampling to equalize class counts", 1, WHITE),
    ("Inverse-frequency class weights for loss function", 1, WHITE),
    ("Fixed seed (42) for full reproducibility", 1, WHITE),
], 19)
bullets(s, 7.0, 1.3, 5.5, 5.5, [
    ("81 Classes in 10 Categories", 0, ACCENT, True),
    ("Animals (21): duck, lion, spider, swan, tiger ...", 1, LGRAY),
    ("Objects (27): laptop, axe, saxophone, sword ...", 1, LGRAY),
    ("Vehicles (5): ladder, bicycle, firetruck, scissors", 1, LGRAY),
    ("Architecture (5): barn, lighthouse, skyscraper", 1, LGRAY),
    ("Sports (4): baseball, basketball, soccer_ball", 1, LGRAY),
    ("Landmarks (2): Eiffel Tower, Mona Lisa", 1, LGRAY),
    ("", 0, None),
    ("Why Real vs. Sketch?", 0, ACCENT, True),
    ("Maximum visual contrast: texture-rich vs. line-only", 1, WHITE),
    ("Tests whether models learn shape vs. texture", 1, WHITE),
    ("Mirrors real deployment gaps (e.g. photo \u2192 X-ray)", 1, WHITE),
], 19)
snum(s, 3)

# ================================================================
# SLIDE 3b: Sample Data — Real vs Sketch
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
title(s, "Sample Data: Real Photos vs. Sketch Drawings")
tb(s, 0.8, 1.2, 11.7, 0.35, "Same object class, radically different visual representation — this is the domain gap our models must bridge.", 16, LGRAY, False, PP_ALIGN.CENTER)

# 6 columns: 3 classes x 2 (real + sketch side by side)
sample_classes = [
    ("Tiger", "tiger", "real_312_000001.jpg", "sketch_312_000001.jpg"),
    ("Bicycle", "bicycle", "real_030_000001.jpg", "sketch_030_000001.jpg"),
    ("Banana", "banana", "real_014_000001.jpg", "sketch_014_000001.jpg"),
    ("Lighthouse", "lighthouse", "real_172_000001.jpg", "sketch_172_000001.jpg"),
    ("Lion", "lion", "real_175_000001.jpg", "sketch_175_000001.jpg"),
    ("Ladder", "ladder", "real_165_000001.jpg", "sketch_165_000001.jpg"),
]

# Layout: 3 columns, 2 rows of classes
col_w = 3.8; col_gap = 0.4
img_w = 1.75; img_h = 1.4
start_x = 0.5

for idx, (pretty, cls, real_fname, sketch_fname) in enumerate(sample_classes):
    col = idx % 3
    row = idx // 3
    x = start_x + col * (col_w + col_gap)
    y_base = 1.65 + row * 3.0

    # Class name
    tb(s, x, y_base, col_w, 0.3, pretty, 16, ACCENT, True, PP_ALIGN.CENTER)

    # Real image
    real_path = f"{BASE}/real/{cls}/{real_fname}"
    sketch_path = f"{BASE}/sketch/{cls}/{sketch_fname}"
    if os.path.exists(real_path):
        img(s, real_path, x, y_base + 0.35, w=img_w, h=img_h)
    tb(s, x, y_base + 0.35 + img_h, img_w, 0.25, "Real", 12, GREEN, True, PP_ALIGN.CENTER)

    # Sketch image
    if os.path.exists(sketch_path):
        img(s, sketch_path, x + img_w + 0.15, y_base + 0.35, w=img_w, h=img_h)
    tb(s, x + img_w + 0.15, y_base + 0.35 + img_h, img_w, 0.25, "Sketch", 12, ORANGE, True, PP_ALIGN.CENTER)

snum(s, 4)

# ================================================================
# SLIDE 4: Model Architecture
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Model Architecture & Training Strategy")
bullets(s, 0.8, 1.3, 5.8, 5.5, [
    ("Architecture: ResNet-152", 0, ORANGE, True),
    ("152-layer residual network, pretrained ImageNet (V2)", 1, WHITE),
    ("Backbone: 2048-D feature extractor", 1, WHITE),
    ("Custom FC head: Linear(2048 \u2192 81 classes)", 1, WHITE),
    ("Multi-GPU training with nn.DataParallel", 1, WHITE),
    ("", 0, None),
    ("Training Techniques", 0, ORANGE, True),
    ("Frozen warm-up (epochs 0-4): freeze conv1, bn1, layer1-2", 1, WHITE),
    ("  \u2192 Prevents destroying pretrained low-level features", 2, LGRAY),
    ("Differential LR: backbone 0.1\u00d7 vs. head 1\u00d7", 1, WHITE),
    ("Cosine annealing LR with linear warmup", 1, WHITE),
    ("Mixup augmentation (\u03b1 = 0.2)", 1, WHITE),
    ("Label smoothing (\u03b5 = 0.1)", 1, WHITE),
    ("Weighted CrossEntropy for class imbalance", 1, WHITE),
    ("Early stopping on macro F1 (patience = 10)", 1, WHITE),
], 18)
bullets(s, 7.0, 1.3, 5.5, 5.5, [
    ("Data Augmentation", 0, ACCENT, True),
    ("Training:", 1, WHITE),
    ("  RandomResizedCrop (0.7-1.0)", 2, LGRAY),
    ("  Random horizontal flip", 2, LGRAY),
    ("  Random rotation (\u00b115\u00b0)", 2, LGRAY),
    ("  Color jitter + random grayscale (10%)", 2, LGRAY),
    ("  Random erasing (20%)", 2, LGRAY),
    ("  ImageNet normalization", 2, LGRAY),
    ("", 0, None),
    ("Validation/Test:", 1, WHITE),
    ("  Resize(256) \u2192 CenterCrop(224) \u2192 Normalize", 2, LGRAY),
    ("", 0, None),
    ("Why These Choices?", 0, ACCENT, True),
    ("Frozen warm-up preserves edge/texture detectors", 1, WHITE),
    ("Mixup blends pairs \u2192 smoother decision boundaries", 1, WHITE),
    ("Differential LR: adapt head fast, update backbone slowly", 1, WHITE),
], 18)
snum(s, 4)

# ================================================================
# SLIDE 5: Training Curves
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Training Curves & Convergence")
tb(s, 0.8, 1.2, 5.8, 0.4, "Real Model Training", 24, ORANGE, True, PP_ALIGN.CENTER)
img(s, f"{EVAL}/real_model_on_real/training_summary.png", 0.3, 1.7, w=6.2)
tb(s, 7.0, 1.2, 5.8, 0.4, "Sketch Model Training", 24, ORANGE, True, PP_ALIGN.CENTER)
img(s, f"{EVAL}/sketch_model_on_sketch/training_summary.png", 6.5, 1.7, w=6.2)
c = card(s, 0.5, 5.8, 12.3, 1.3, ACCENT)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.text = "Observations"; p.font.size = Pt(18); p.font.color.rgb = ACCENT; p.font.bold = True
for t in ["Real model: fast convergence \u2014 88.9% val acc at epoch 1 \u2192 94.6% by epoch 9. Best val F1 = 0.9451 (epoch 9)",
          "Sketch model: slower convergence \u2014 58.2% at epoch 1 \u2192 83.1% by epoch 13. Best val F1 = 0.8336. Harder task: sparse line drawings lack texture cues"]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(15); p.font.color.rgb = WHITE
snum(s, 5)

# ================================================================
# SLIDE 6: Classification Results
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Classification Performance")
for col, txt, w in [(1.0, "Scenario", 4.5), (5.8, "Accuracy", 1.5), (7.5, "F1 (macro)", 1.5), (9.2, "Precision", 1.5), (10.9, "Recall", 1.5)]:
    tb(s, col, 1.3, w, 0.4, txt, 17, ACCENT, True, PP_ALIGN.CENTER)
for i, (sc, acc, f1, pr, rc, cl) in enumerate([
    ("Real \u2192 Real (in-domain)", "93.63%", "0.9288", "0.9276", "0.9334", GREEN),
    ("Sketch \u2192 Sketch (in-domain)", "80.69%", "0.8096", "0.8111", "0.8187", GREEN),
    ("Real \u2192 Sketch (cross-domain)", "56.13%", "0.5692", "0.6803", "0.5751", RED),
    ("Sketch \u2192 Real (cross-domain)", "82.99%", "0.8161", "0.8423", "0.8191", GREEN),
]):
    y = 1.8 + i * 0.5
    tb(s, 1.0, y, 4.5, 0.4, sc, 17, WHITE)
    for cx, v in [(5.8, acc), (7.5, f1), (9.2, pr), (10.9, rc)]:
        tb(s, cx, y, 1.5, 0.4, v, 17, cl, False, PP_ALIGN.CENTER)
# Confusion matrices
for j, (lbl, path, clr) in enumerate([
    ("Real \u2192 Real", f"{DOCS}/evaluation/cm_real_real.png", ACCENT),
    ("Real \u2192 Sketch", f"{DOCS}/evaluation/cm_real_sketch.png", RED),
    ("Sketch \u2192 Sketch", f"{DOCS}/evaluation/cm_sketch_sketch.png", ACCENT),
    ("Sketch \u2192 Real", f"{DOCS}/evaluation/cm_sketch_real.png", GREEN),
]):
    x = 0.3 + j * 3.2
    tb(s, x, 4.0, 3.0, 0.35, lbl, 15, clr, True, PP_ALIGN.CENTER)
    img(s, path, x, 4.35, w=3.0)
snum(s, 6)

# ================================================================
# SLIDE 7: Asymmetric Transfer
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Key Finding: Asymmetric Domain Transfer")
c = card(s, 0.5, 1.3, 5.8, 5.8, RED)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "Real Model \u2192 Sketch: 37.5-point DROP"; p.font.size = Pt(22); p.font.color.rgb = RED; p.font.bold = True
for t, c2 in [
    ("In-domain: 93.63%  \u2192  Cross-domain: 56.13%", WHITE),
    ("", WHITE), ("Why it fails:", YELLOW),
    ("\u2022 Photos are rich in texture, color, lighting", LGRAY),
    ("\u2022 CNN learns texture-dependent representations", LGRAY),
    ("\u2022 Sketches have NO texture \u2014 only lines & shapes", LGRAY),
    ("\u2022 Features the model relies on don't exist in sketches", LGRAY),
    ("", WHITE), ("Evidence from t-SNE:", YELLOW),
    ("\u2022 Real & sketch images form SEPARATE clusters", LGRAY),
    ("\u2022 Same-class images from different domains are far apart", LGRAY),
    ("", WHITE), ("This is the 'texture bias' problem (Geirhos 2019):", YELLOW),
    ("CNNs are biased toward texture over shape.", WHITE),
]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(16); p.font.color.rgb = c2

c = card(s, 6.8, 1.3, 5.8, 5.8, GREEN)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "Sketch Model \u2192 Real: only 0.3-point drop!"; p.font.size = Pt(22); p.font.color.rgb = GREEN; p.font.bold = True
for t, c2 in [
    ("In-domain: 80.69%  \u2192  Cross-domain: 82.99%", WHITE),
    ("", WHITE), ("Why it succeeds:", YELLOW),
    ("\u2022 Sketches force model to learn from SHAPE only", LGRAY),
    ("\u2022 Shape-based features are universal across domains", LGRAY),
    ("\u2022 Real photos contain shape + extra texture info", LGRAY),
    ("\u2022 Extra information helps \u2192 slight accuracy gain", LGRAY),
    ("", WHITE), ("Evidence from t-SNE:", YELLOW),
    ("\u2022 Real & sketch images are INTERLEAVED in feature space", LGRAY),
    ("\u2022 Same-class images from both domains overlap", LGRAY),
    ("", WHITE), ("Practical implication:", YELLOW),
    ("Training on abstract/minimal data can produce", WHITE),
    ("more robust, generalizable models.", WHITE),
]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(16); p.font.color.rgb = c2
snum(s, 7)

# ================================================================
# SLIDE 8: XAI Methods Theory
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "XAI Methods: How Each One Works")
method_info = [
    ("Grad-CAM", ORANGE, [
        "How: Compute gradient of class score",
        "w.r.t. final conv layer feature maps.",
        "Weight each channel by its gradient,",
        "then average \u2192 activation map.",
        "",
        "Output: 7\u00d77 coarse heatmap",
        "(upsampled to 224\u00d7224)",
        "Speed: 1 fwd + 1 bwd pass (fastest)",
        "",
        "Intuition: Which spatial regions in",
        "the final layer activate most for",
        "this class?",
        "",
        "Limitation: Coarse resolution \u2014",
        "highlights regions, not pixels.",
    ]),
    ("Grad-CAM++", GREEN, [
        "How: Like Grad-CAM but uses",
        "higher-order partial derivatives",
        "(2nd & 3rd order) to compute",
        "pixel-wise weights.",
        "",
        "Output: 7\u00d77 heatmap",
        "(better full-object coverage)",
        "Speed: Similar to Grad-CAM (fast)",
        "",
        "Advantage: Better coverage of the",
        "entire object, not just the single",
        "most discriminative sub-region.",
        "",
        "Better for multi-instance images",
        "and images with occluded objects.",
    ]),
    ("Integrated Gradients", ACCENT, [
        "How: Interpolate from a black",
        "baseline to the actual input in",
        "101 steps. Compute gradient at",
        "each step, then average.",
        "",
        "Output: 224\u00d7224 pixel-level map",
        "Speed: 101 fwd+bwd (~100\u00d7 slower)",
        "",
        "Axioms satisfied:",
        "\u2022 Sensitivity: important features",
        "  get non-zero attribution",
        "\u2022 Completeness: attributions sum",
        "  to output \u2212 baseline difference",
        "",
        "Limitation: Noisy, baseline-dependent.",
    ]),
    ("LIME", PURPLE, [
        "How: Segment image into super-",
        "pixels. Generate 3,000 perturbed",
        "versions (randomly mask super-",
        "pixels). Query model on each.",
        "Fit linear model to predict output.",
        "",
        "Output: Per-superpixel weights",
        "Speed: ~3,000 fwd passes (slowest)",
        "",
        "Key: Completely model-agnostic \u2014",
        "treats network as a black box.",
        "Most human-interpretable output.",
        "",
        "Limitation: Stochastic \u2014 different",
        "runs give slightly different results.",
    ]),
]
for i, (name, color, lines) in enumerate(method_info):
    left = 0.4 + i * 3.2
    c = card(s, left, 1.2, 3.05, 6.0, color)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(21); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    for line in lines:
        p = tf.add_paragraph(); p.text = line; p.font.size = Pt(13)
        p.font.color.rgb = WHITE if not line.startswith(" ") and not line.startswith("\u2022") else LGRAY
        p.space_after = Pt(1)
snum(s, 8)

# ================================================================
# SLIDE 9: Domain Shift Introduction
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "The Challenge: Domain Shift")
c = card(s, 0.5, 1.3, 12.3, 2.5, ACCENT)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "What is Domain Shift?"; p.font.size = Pt(24); p.font.color.rgb = ACCENT; p.font.bold = True
for t, clr in [
    ("A critical challenge in deploying deep learning models \u2014 the performance degradation that occurs when a model trained on one data distribution is applied to a different but related distribution.", WHITE),
    ("Domain shift is pervasive: a medical imaging model trained at one hospital may fail at another; an autonomous driving system trained in sunny conditions may struggle in rain or fog.", LGRAY),
    ("Understanding HOW a model's reasoning changes under domain shift is arguably even more important than measuring the accuracy drop itself.", YELLOW),
    ("It reveals whether the model has learned genuinely transferable concepts or merely domain-specific shortcuts.", ORANGE),
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(16); p.font.color.rgb = clr; p.space_after = Pt(4)

c2 = card(s, 0.5, 4.1, 12.3, 3.0, ORANGE)
tf2 = c2.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0.25); tf2.margin_top = Inches(0.12)
p2 = tf2.paragraphs[0]; p2.text = "Why This Matters for Explainability"; p2.font.size = Pt(22); p2.font.color.rgb = ORANGE; p2.font.bold = True
for t, clr in [
    ("Even if a model maintains high accuracy across domains, its EXPLANATIONS may shift dramatically", WHITE),
    ("A correct prediction with an incorrect explanation is dangerous in safety-critical applications", YELLOW),
    ("If explanations change under domain shift, we cannot trust them for clinical decisions, autonomous driving, or legal compliance", WHITE),
    ("Our study systematically measures explanation reliability across visual domains using 4 metrics", LGRAY),
]:
    p2 = tf2.add_paragraph(); p2.text = "\u2022 " + t; p2.font.size = Pt(15); p2.font.color.rgb = clr; p2.space_after = Pt(4)
snum(s, 9)

# ================================================================
# SLIDE 10: Top & Bottom Classes — Real Model
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Top & Bottom Classes: Real Model (Real \u2192 Real)")

tb(s, 0.4, 1.25, 6.2, 0.4, "Top 10 Classes", 20, GREEN, True, PP_ALIGN.CENTER)
for hdr, x, w in [("Class", 0.4, 2.8), ("Acc%", 3.3, 1.3), ("F1", 4.7, 1.3)]:
    tb(s, x, 1.7, w, 0.35, hdr, 14, ACCENT, True, PP_ALIGN.CENTER)
top10_real = [
    ("The Mona Lisa", "100.0", "0.983"), ("basketball", "100.0", "1.000"), ("saxophone", "100.0", "1.000"),
    ("skull", "100.0", "1.000"), ("trumpet", "100.0", "1.000"), ("backpack", "100.0", "0.936"),
    ("banana", "100.0", "0.963"), ("baseball", "100.0", "0.900"), ("firetruck", "100.0", "0.974"),
    ("spreadsheet", "100.0", "0.987"),
]
for j, (cls, acc, f1) in enumerate(top10_real):
    y = 2.1 + j * 0.38
    tb(s, 0.4, y, 2.8, 0.35, cls, 13, WHITE); tb(s, 3.3, y, 1.3, 0.35, acc, 13, GREEN, False, PP_ALIGN.CENTER); tb(s, 4.7, y, 1.3, 0.35, f1, 13, GREEN, False, PP_ALIGN.CENTER)

tb(s, 6.8, 1.25, 6.2, 0.4, "Bottom 10 Classes", 20, RED, True, PP_ALIGN.CENTER)
for hdr, x, w in [("Class", 6.8, 2.8), ("Acc%", 9.7, 1.3), ("F1", 11.1, 1.3)]:
    tb(s, x, 1.7, w, 0.35, hdr, 14, ACCENT, True, PP_ALIGN.CENTER)
bot10_real = [
    ("saw", "66.67", "0.696"), ("eye", "75.36", "0.846"), ("circle", "80.77", "0.857"),
    ("spoon", "81.13", "0.878"), ("triangle", "81.58", "0.838"), ("flower", "83.33", "0.896"),
    ("barn", "83.87", "0.881"), ("shovel", "84.44", "0.905"), ("feather", "86.00", "0.860"),
    ("suitcase", "86.05", "0.914"),
]
for j, (cls, acc, f1) in enumerate(bot10_real):
    y = 2.1 + j * 0.38
    tb(s, 6.8, y, 2.8, 0.35, cls, 13, WHITE); tb(s, 9.7, y, 1.3, 0.35, acc, 13, RED, False, PP_ALIGN.CENTER); tb(s, 11.1, y, 1.3, 0.35, f1, 13, RED, False, PP_ALIGN.CENTER)

c2 = card(s, 0.4, 5.95, 12.5, 1.2, ORANGE)
tf2 = c2.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0.15); tf2.margin_top = Inches(0.06)
p2 = tf2.paragraphs[0]; p2.text = "Analysis"; p2.font.size = Pt(16); p2.font.color.rgb = ORANGE; p2.font.bold = True
p2 = tf2.add_paragraph(); p2.text = "Top performers have highly distinctive visual features (saxophone's curved metallic shape, trumpet's bell). Bottom performers: saw (66.67%) confused with elongated tools; circle and triangle lack rich texture cues; spoon shares similarity with other utensils."; p2.font.size = Pt(13); p2.font.color.rgb = WHITE
snum(s, 10)

# ================================================================
# SLIDE 11: Top & Bottom Classes — Sketch Model
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Top & Bottom Classes: Sketch Model (Sketch \u2192 Sketch)")

tb(s, 0.4, 1.25, 6.2, 0.4, "Top 10 Classes", 20, GREEN, True, PP_ALIGN.CENTER)
for hdr, x, w in [("Class", 0.4, 2.8), ("Acc%", 3.3, 1.3), ("F1", 4.7, 1.3)]:
    tb(s, x, 1.7, w, 0.35, hdr, 14, ACCENT, True, PP_ALIGN.CENTER)
top10_sketch = [
    ("banana", "100.0", "0.952"), ("barn", "100.0", "0.930"), ("baseball", "100.0", "0.833"),
    ("bicycle", "100.0", "1.000"), ("marker", "100.0", "0.828"), ("axe", "95.45", "0.840"),
    ("birthday cake", "95.65", "0.936"), ("light bulb", "95.00", "0.905"),
    ("tiger", "94.87", "0.925"), ("laptop", "93.75", "0.896"),
]
for j, (cls, acc, f1) in enumerate(top10_sketch):
    y = 2.1 + j * 0.38
    tb(s, 0.4, y, 2.8, 0.35, cls, 13, WHITE); tb(s, 3.3, y, 1.3, 0.35, acc, 13, GREEN, False, PP_ALIGN.CENTER); tb(s, 4.7, y, 1.3, 0.35, f1, 13, GREEN, False, PP_ALIGN.CENTER)

tb(s, 6.8, 1.25, 6.2, 0.4, "Bottom 10 Classes", 20, RED, True, PP_ALIGN.CENTER)
for hdr, x, w in [("Class", 6.8, 2.8), ("Acc%", 9.7, 1.3), ("F1", 11.1, 1.3)]:
    tb(s, x, 1.7, w, 0.35, hdr, 14, ACCENT, True, PP_ALIGN.CENTER)
bot10_sketch = [
    ("shovel", "44.44", "0.544"), ("scorpion", "55.56", "0.694"), ("sword", "55.26", "0.568"),
    ("firetruck", "57.58", "0.623"), ("triangle", "60.00", "0.667"), ("saw", "63.64", "0.667"),
    ("palm tree", "64.71", "0.550"), ("sock", "64.44", "0.699"),
    ("circle", "65.00", "0.578"), ("scissors", "65.91", "0.674"),
]
for j, (cls, acc, f1) in enumerate(bot10_sketch):
    y = 2.1 + j * 0.38
    tb(s, 6.8, y, 2.8, 0.35, cls, 13, WHITE); tb(s, 9.7, y, 1.3, 0.35, acc, 13, RED, False, PP_ALIGN.CENTER); tb(s, 11.1, y, 1.3, 0.35, f1, 13, RED, False, PP_ALIGN.CENTER)

c2 = card(s, 0.4, 5.95, 12.5, 1.2, ORANGE)
tf2 = c2.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0.15); tf2.margin_top = Inches(0.06)
p2 = tf2.paragraphs[0]; p2.text = "Analysis"; p2.font.size = Pt(16); p2.font.color.rgb = ORANGE; p2.font.bold = True
p2 = tf2.add_paragraph(); p2.text = "Top performers: structurally distinctive shapes (bicycle's two-wheel frame, banana's curve). Bottom: elongated tools hardest (shovel 44.44%, sword 55.26%). Firetruck drops 100% (real) \u2192 57.58% (sketch) \u2014 sketches lack distinctive red color and detailed features."; p2.font.size = Pt(13); p2.font.color.rgb = WHITE
snum(s, 11)

# ================================================================
# SLIDES 12+: Attribution Map Image Slides (2x2 grid, height-constrained)
# ================================================================
slide_counter = [12]

# --- 1. Real Model on Real Domain (in-domain) ---
for cls in ['ladder', 'lion', 'tiger', 'lighthouse', 'banana', 'spider', 'snake', 'umbrella', 'airplane']:
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    title(s, "Real Model \u2014 Real Domain (In-Domain)")
    pretty = cls.replace('_', ' ').title()
    tb(s, 0.5, 1.2, 12.3, 0.35, f"Class: {pretty}", 22, YELLOW, True, PP_ALIGN.CENTER)
    add_2x2_grid(s, 'real_model', 'real', cls, y_start=1.6)
    snum(s, slide_counter[0]); slide_counter[0] += 1

# --- 2. Sketch Model on Real Domain (cross-domain — the success case) ---
for cls in ['ladder', 'tiger', 'lion', 'banana', 'lighthouse', 'bicycle', 'snake', 'umbrella', 'airplane']:
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    title(s, "Sketch Model \u2014 Real Domain (Cross-Domain)")
    pretty = cls.replace('_', ' ').title()
    tb(s, 0.5, 1.2, 12.3, 0.35, f"Class: {pretty}", 22, YELLOW, True, PP_ALIGN.CENTER)
    tb(s, 0.5, 1.5, 12.3, 0.25, "Sketch-trained model explains real photos \u2014 the successful transfer case!", 14, LGRAY, False, PP_ALIGN.CENTER)
    add_2x2_grid(s, 'sketch_model', 'real', cls, y_start=1.8)
    snum(s, slide_counter[0]); slide_counter[0] += 1

# --- 3. Real Model on Sketch Domain (cross-domain) ---
for cls in ['ladder', 'tiger', 'bicycle', 'lion', 'snake', 'umbrella', 'airplane']:
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    title(s, "Real Model \u2014 Sketch Domain (Cross-Domain)")
    pretty = cls.replace('_', ' ').title()
    tb(s, 0.5, 1.2, 12.3, 0.35, f"Class: {pretty}", 22, YELLOW, True, PP_ALIGN.CENTER)
    tb(s, 0.5, 1.5, 12.3, 0.25, "How does the Real-trained model explain sketch images it was never trained on?", 14, LGRAY, False, PP_ALIGN.CENTER)
    add_2x2_grid(s, 'real_model', 'sketch', cls, y_start=1.8)
    snum(s, slide_counter[0]); slide_counter[0] += 1

# --- 4. Sketch Model on Sketch Domain (in-domain) ---
for cls in ['ladder', 'tiger', 'bicycle', 'saxophone', 'lion', 'scissors', 'snake', 'umbrella', 'airplane']:
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    title(s, "Sketch Model \u2014 Sketch Domain (In-Domain)")
    pretty = cls.replace('_', ' ').title()
    tb(s, 0.5, 1.2, 12.3, 0.35, f"Class: {pretty}", 22, YELLOW, True, PP_ALIGN.CENTER)
    add_2x2_grid(s, 'sketch_model', 'sketch', cls, y_start=1.6)
    snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# FAILURE CASE SLIDE 1: Shovel & Scorpion
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
title(s, "Failure Analysis (1/2): Shovel & Scorpion")

# --- Shovel ---
tb(s, 0.4, 1.2, 12.5, 0.35, "Shovel \u2014 Worst Sketch Model performer (44.44% accuracy)", 18, RED, True)
fail_iw = 4.0; fail_ih = 1.45
tb(s, 0.5, 1.6, fail_iw, 0.25, "Grad-CAM (Real Model on Sketch)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/gradcam/sketch/shovel_0961.png", 0.5, 1.87, w=fail_iw, h=fail_ih)
tb(s, 0.5 + fail_iw + 0.2, 1.6, fail_iw, 0.25, "LIME (Real Model on Sketch)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/lime/sketch/shovel_0961.png", 0.5 + fail_iw + 0.2, 1.87, w=fail_iw, h=fail_ih)

c1 = card(s, 8.9, 1.55, 4.1, 1.85, RED)
tf1 = c1.text_frame; tf1.word_wrap = True; tf1.margin_left = Inches(0.1); tf1.margin_top = Inches(0.05)
p1 = tf1.paragraphs[0]; p1.text = "Why it fails:"; p1.font.size = Pt(14); p1.font.color.rgb = RED; p1.font.bold = True
for txt in [
    "Sketch shows a person digging \u2014 model focuses on the PERSON, not the shovel",
    "Shovel is a secondary object in the scene, visually overshadowed",
    "Elongated shape confused with axe, sword, screwdriver",
]:
    p1 = tf1.add_paragraph()
    p1.text = "\u2022 " + txt; p1.font.size = Pt(11); p1.font.color.rgb = WHITE; p1.space_after = Pt(2)

# --- Scorpion ---
tb(s, 0.4, 3.65, 12.5, 0.35, "Scorpion \u2014 55.56% accuracy on Sketch Model", 18, RED, True)
tb(s, 0.5, 4.05, fail_iw, 0.25, "Grad-CAM (Real Model on Sketch)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/gradcam/sketch/scorpion_0861.png", 0.5, 4.32, w=fail_iw, h=fail_ih)
tb(s, 0.5 + fail_iw + 0.2, 4.05, fail_iw, 0.25, "LIME (Real Model on Sketch)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/lime/sketch/scorpion_0861.png", 0.5 + fail_iw + 0.2, 4.32, w=fail_iw, h=fail_ih)

c2 = card(s, 8.9, 4.0, 4.1, 1.85, RED)
tf2 = c2.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0.1); tf2.margin_top = Inches(0.05)
p2 = tf2.paragraphs[0]; p2.text = "Why it fails:"; p2.font.size = Pt(14); p2.font.color.rgb = RED; p2.font.bold = True
for txt in [
    "Detailed sketch with many legs \u2014 model attention is diffuse across the body",
    "Scorpion confused with spider and crab (similar multi-legged arthropod shape)",
    "Texture cues (exoskeleton sheen) absent in sketch \u2014 real model has no shape-only reference",
]:
    p2 = tf2.add_paragraph()
    p2.text = "\u2022 " + txt; p2.font.size = Pt(11); p2.font.color.rgb = WHITE; p2.space_after = Pt(2)

# Bottom takeaway
c3 = card(s, 0.3, 6.1, 12.7, 1.1, YELLOW)
tf3 = c3.text_frame; tf3.word_wrap = True; tf3.margin_left = Inches(0.15); tf3.margin_top = Inches(0.05)
p3 = tf3.paragraphs[0]; p3.text = "Key Insight"; p3.font.size = Pt(15); p3.font.color.rgb = YELLOW; p3.font.bold = True
p3 = tf3.add_paragraph()
p3.text = "When the model misclassifies, the attribution maps reveal WHY: the model attends to wrong objects (shovel \u2192 person) or distributes attention diffusely (scorpion). Poor predictions produce poor explanations \u2014 explanation quality is tightly coupled to model confidence."
p3.font.size = Pt(13); p3.font.color.rgb = WHITE
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# FAILURE CASE SLIDE 2: Scissors (cross-domain)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
title(s, "Failure Analysis (2/2): Scissors Under Domain Shift")

tb(s, 0.4, 1.2, 12.5, 0.35, "Scissors \u2014 Real model predicts sketch scissors as 'baseball_bat'", 18, RED, True)

# Row 1: Real model on sketch scissors
tb(s, 0.5, 1.65, fail_iw, 0.25, "Grad-CAM (Real Model on Sketch)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/gradcam/sketch/scissors_0846.png", 0.5, 1.92, w=fail_iw, h=fail_ih)
tb(s, 0.5 + fail_iw + 0.2, 1.65, fail_iw, 0.25, "LIME (Real Model on Sketch)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/lime/sketch/scissors_0846.png", 0.5 + fail_iw + 0.2, 1.92, w=fail_iw, h=fail_ih)

c4 = card(s, 8.9, 1.6, 4.1, 1.85, RED)
tf4 = c4.text_frame; tf4.word_wrap = True; tf4.margin_left = Inches(0.1); tf4.margin_top = Inches(0.05)
p4 = tf4.paragraphs[0]; p4.text = "Why it fails:"; p4.font.size = Pt(14); p4.font.color.rgb = RED; p4.font.bold = True
for txt in [
    "Sketch scissors drawn as pliers/tongs \u2014 visually dissimilar to photo scissors",
    "Model focuses on the handle junction, sees elongated shape \u2192 predicts 'baseball_bat'",
    "Real model relies on texture (metallic blade) which is absent in the sketch",
]:
    p4 = tf4.add_paragraph()
    p4.text = "\u2022 " + txt; p4.font.size = Pt(11); p4.font.color.rgb = WHITE; p4.space_after = Pt(2)

# Row 2: Contrast with real model on REAL scissors (success case)
tb(s, 0.4, 3.65, 12.5, 0.35, "Compare: Same class, real domain \u2014 correctly classified (98.25% accuracy)", 18, GREEN, True)
tb(s, 0.5, 4.05, fail_iw, 0.25, "Grad-CAM (Real Model on Real)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/gradcam/real/scissors_0865.png", 0.5, 4.32, w=fail_iw, h=fail_ih)
tb(s, 0.5 + fail_iw + 0.2, 4.05, fail_iw, 0.25, "LIME (Real Model on Real)", 13, ACCENT, True, PP_ALIGN.CENTER)
img(s, f"{XAI_OUT}/real_model/attribution_maps/lime/real/scissors_0865.png", 0.5 + fail_iw + 0.2, 4.32, w=fail_iw, h=fail_ih)

c5 = card(s, 8.9, 4.0, 4.1, 1.85, GREEN)
tf5 = c5.text_frame; tf5.word_wrap = True; tf5.margin_left = Inches(0.1); tf5.margin_top = Inches(0.05)
p5 = tf5.paragraphs[0]; p5.text = "Why it succeeds:"; p5.font.size = Pt(14); p5.font.color.rgb = GREEN; p5.font.bold = True
for txt in [
    "Real photo has texture cues: metallic blades, reflections, colour contrast",
    "Model confidently focuses on blade intersection and handle",
    "Attribution maps are sharp and focused \u2014 high-quality explanations",
]:
    p5 = tf5.add_paragraph()
    p5.text = "\u2022 " + txt; p5.font.size = Pt(11); p5.font.color.rgb = WHITE; p5.space_after = Pt(2)

# Bottom takeaway
c6 = card(s, 0.3, 6.1, 12.7, 1.1, YELLOW)
tf6 = c6.text_frame; tf6.word_wrap = True; tf6.margin_left = Inches(0.15); tf6.margin_top = Inches(0.05)
p6 = tf6.paragraphs[0]; p6.text = "Domain Shift Impact"; p6.font.size = Pt(15); p6.font.color.rgb = YELLOW; p6.font.bold = True
p6 = tf6.add_paragraph()
p6.text = "The same model on the same class shows dramatically different explanation quality depending on the domain. On real photos, attributions are focused and correct. On sketches, the model attends to wrong features and misclassifies \u2014 demonstrating that domain shift degrades both accuracy AND explanation quality."
p6.font.size = Pt(13); p6.font.color.rgb = WHITE
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# NEW: Model Comparison Slides — Real vs Sketch Model, same image
# For tiger, lion, bicycle on both real and sketch images
# Each class+domain = 2 slides (2 methods per slide for larger images)
# ================================================================
for cls_name in ['tiger', 'lion', 'bicycle']:
    add_model_comparison(slide_counter, cls_name, 'real', 'Real')
    add_model_comparison(slide_counter, cls_name, 'sketch', 'Sketch')

# ================================================================
# Evaluation Framework — SLIDE A: Stability + Deletion (2 metrics)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Evaluation Metrics Explained (1/2)")
tb(s, 0.5, 1.15, 12.3, 0.3, "How we measure whether explanations are trustworthy and reliable", 17, LGRAY, False, PP_ALIGN.CENTER)

# Metric 1: Stability
c = card(s, 0.4, 1.6, 12.5, 2.5, GREEN)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.text = "1. Stability (SSIM)                                          \u2191 HIGHER = BETTER"; p.font.size = Pt(20); p.font.color.rgb = GREEN; p.font.bold = True
for t, clr in [
    ("Question: If we add tiny, invisible noise to the image, does the explanation change?", WHITE),
    ("How it works: Add small Gaussian noise to the input image \u2192 recompute the attribution map \u2192 compare original vs. noisy using SSIM (Structural Similarity Index)", LGRAY),
    ("Score range: 0.0 to 1.0    |    1.0 = perfectly identical maps    |    > 0.9 = very stable    |    < 0.5 = unreliable", YELLOW),
    ("Why it matters: If invisible noise completely changes where the model \"looks,\" the explanation cannot be trusted in practice. A doctor needs the same explanation every time for the same image.", WHITE),
    ("Intuition: Think of it like asking someone the same question twice \u2014 a reliable expert gives the same answer. An unreliable one changes their mind with tiny irrelevant changes.", LGRAY),
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = clr; p.space_after = Pt(3)

# Metric 2: Deletion
c = card(s, 0.4, 4.3, 12.5, 2.8, RED)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.text = "2. Faithfulness \u2014 Deletion AUC                      \u2193 LOWER = BETTER"; p.font.size = Pt(20); p.font.color.rgb = RED; p.font.bold = True
for t, clr in [
    ("Question: Do the highlighted pixels ACTUALLY drive the model's prediction?", WHITE),
    ("How it works: Rank all pixels by attribution importance \u2192 progressively remove the MOST important pixels first \u2192 measure how fast the model's confidence drops", LGRAY),
    ("Plot: X-axis = % of pixels removed, Y-axis = model confidence. Compute area under this curve (AUC).", LGRAY),
    ("Score interpretation: LOWER AUC = BETTER. If confidence drops fast when you remove highlighted pixels, those pixels truly mattered.", YELLOW),
    ("A HIGH deletion AUC means removing the \"important\" pixels barely hurt confidence \u2192 the explanation was MISLEADING \u2014 it highlighted the wrong things.", ORANGE),
    ("Analogy: Removing the foundation of a building should cause collapse. If removing it does nothing, something is wrong with your blueprint.", WHITE),
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = clr; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Evaluation Framework — SLIDE B: Insertion + Consistency (2 metrics)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Evaluation Metrics Explained (2/2)")
tb(s, 0.5, 1.15, 12.3, 0.3, "How we measure whether explanations are trustworthy and reliable", 17, LGRAY, False, PP_ALIGN.CENTER)

# Metric 3: Insertion
c = card(s, 0.4, 1.6, 12.5, 2.5, ORANGE)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.text = "3. Faithfulness \u2014 Insertion AUC                      \u2191 HIGHER = BETTER"; p.font.size = Pt(20); p.font.color.rgb = ORANGE; p.font.bold = True
for t, clr in [
    ("Question: Can the highlighted pixels ALONE recover the model's prediction?", WHITE),
    ("How it works: Start with a completely blank/blurred image \u2192 progressively reveal the MOST important pixels first \u2192 measure how fast model confidence rises", LGRAY),
    ("Plot: X-axis = % of pixels revealed, Y-axis = model confidence. Compute area under this curve (AUC).", LGRAY),
    ("Score interpretation: HIGHER AUC = BETTER. If confidence rises fast with just the highlighted pixels, the explanation found the truly important features.", YELLOW),
    ("A LOW insertion AUC means showing the \"important\" pixels barely helps \u2192 the method missed the key features the model actually uses.", ORANGE),
    ("Analogy: A good summary of a book should let you understand the plot. If it doesn't, the summary picked the wrong parts.", WHITE),
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = clr; p.space_after = Pt(3)

# Metric 4: Cross-Domain Consistency
c = card(s, 0.4, 4.3, 12.5, 2.8, ACCENT)
tf = c.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.text = "4. Cross-Domain Consistency (Cosine Sim.)     \u2191 HIGHER = BETTER"; p.font.size = Pt(20); p.font.color.rgb = ACCENT; p.font.bold = True
for t, clr in [
    ("Question: Does the model focus on the SAME regions regardless of visual style (photo vs. sketch)?", WHITE),
    ("How it works: For the same class, compute attribution maps on real images AND sketch images \u2192 flatten to vectors \u2192 measure cosine similarity between them", LGRAY),
    ("Score range: -1.0 to 1.0    |    1.0 = identical focus    |    > 0.9 = very consistent    |    < 0.8 = explanations shift significantly", YELLOW),
    ("HIGHER = BETTER. A score of 0.95 means the model looks at nearly the same spatial regions in both photos and sketches.", GREEN),
    ("A LOW score means the explanation changes dramatically depending on visual style \u2192 the model uses different reasoning strategies for different domains.", RED),
    ("Why critical: For safety-critical deployment (medical AI, autonomous driving), explanations MUST be stable across visual conditions and equipment.", WHITE),
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = clr; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Stability Results — Table
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Stability Results (SSIM \u2014 \u2191 Higher is Better)")
tb(s, 0.5, 1.15, 12.3, 0.35, "SSIM compares attribution maps before & after adding Gaussian noise. Range: 0\u21921. Higher = more stable.", 16, LGRAY)

headers = ["Method", "Real Model\nReal Images", "Real Model\nSketch Images", "Sketch Model\nReal Images", "Sketch Model\nSketch Images"]
cx = [0.5, 3.2, 5.3, 7.4, 9.5]
cw = [2.5, 2.1, 2.1, 2.1, 2.1]
for i, (x, h) in enumerate(zip(cx, headers)):
    tb(s, x, 1.6, cw[i], 0.7, h, 15, ACCENT, True, PP_ALIGN.CENTER)
data = [
    ("Grad-CAM++", "0.946 \u00b1 0.038", "0.927 \u00b1 0.050", "0.946 \u00b1 0.042", "0.940 \u00b1 0.043", GREEN),
    ("Grad-CAM",   "0.942 \u00b1 0.056", "0.896 \u00b1 0.094", "0.928 \u00b1 0.078", "0.926 \u00b1 0.079", GREEN),
    ("Integ. Grad.","0.563 \u00b1 0.162", "0.555 \u00b1 0.156", "0.542 \u00b1 0.181", "0.492 \u00b1 0.170", YELLOW),
    ("LIME",       "0.471 \u00b1 0.136", "0.471 \u00b1 0.137", "0.485 \u00b1 0.139", "0.493 \u00b1 0.143", RED),
]
for j, (m, *vals, clr) in enumerate(data):
    y = 2.4 + j * 0.5
    tb(s, cx[0], y, cw[0], 0.45, m, 16, WHITE, True)
    for k, v in enumerate(vals):
        tb(s, cx[k+1], y, cw[k+1], 0.45, v, 15, clr, False, PP_ALIGN.CENTER)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Stability Results — Interpretation (separate slide)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Stability: What Do These Numbers Mean?")

c2 = card(s, 0.4, 1.3, 12.5, 5.8, GREEN)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "Interpretation & Analysis"; p.font.size = Pt(22); p.font.color.rgb = GREEN; p.font.bold = True
for t, c3 in [
    ("Grad-CAM & Grad-CAM++ achieve SSIM > 0.90 in ALL conditions \u2014 highly stable explanations", WHITE),
    ("Why? Their 7\u00d77 spatial resolution acts as a natural low-pass filter. Small input perturbations are smoothed away by the coarse grid.", LGRAY),
    ("", WHITE),
    ("Integrated Gradients (SSIM ~0.55): operates at full pixel level (224\u00d7224). Noise directly shifts individual pixel attributions \u2192 inherently less stable.", WHITE),
    ("", WHITE),
    ("LIME (SSIM ~0.47): stochastic by design \u2014 each run samples 3,000 random perturbations. Different noise \u2192 different samples \u2192 different explanations.", WHITE),
    ("", WHITE),
    ("Fundamental trade-off: higher-resolution explanations are LESS stable. You get either coarse-but-reliable OR fine-grained-but-noisy.", YELLOW),
    ("", WHITE),
    ("Cross-domain effect: Real model's stability drops on sketches (0.946\u21920.927 for GradCAM++) \u2014 domain shift reduces explanation consistency.", YELLOW),
]:
    p = tf.add_paragraph(); p.text = ("\u2022 " + t) if t else ""; p.font.size = Pt(15); p.font.color.rgb = c3; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Faithfulness — Deletion Table
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Faithfulness: Deletion AUC (\u2193 Lower is Better)")

tb(s, 0.5, 1.15, 12.3, 0.5, "Remove the most important pixels first \u2192 model confidence should DROP.\nLower AUC = explanation correctly identified important regions.", 16, LGRAY)

dh = ["Method", "Real Model\nReal Images", "Real Model\nSketch Images", "Sketch Model\nReal Images", "Sketch Model\nSketch Images"]
dx = [0.8, 3.5, 5.6, 7.7, 9.8]
dw = [2.5, 2.1, 2.1, 2.1, 2.1]
for i, (x, h) in enumerate(zip(dx, dh)):
    tb(s, x, 1.7, dw[i], 0.7, h, 15, ACCENT, True, PP_ALIGN.CENTER)

del_data = [
    ("Grad-CAM++", ["0.323", "0.125", "0.205", "0.194"]),
    ("Grad-CAM",   ["0.294", "0.107", "0.182", "0.170"]),
    ("Integ. Grad.",["0.277", "0.125", "0.177", "0.172"]),
    ("LIME",       ["0.152", "0.054", "0.071", "0.067"]),
]
for j, (m, vals) in enumerate(del_data):
    y = 2.5 + j * 0.55
    tb(s, dx[0], y, dw[0], 0.45, m, 16, WHITE, True)
    for k, v in enumerate(vals):
        fv = float(v)
        clr = GREEN if fv < 0.10 else (YELLOW if fv < 0.15 else (ORANGE if fv < 0.25 else WHITE))
        tb(s, dx[k+1], y, dw[k+1], 0.45, v, 16, clr, fv < 0.10, PP_ALIGN.CENTER)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Deletion — Interpretation (separate slide)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Deletion: What Do These Numbers Mean?")

c2 = card(s, 0.4, 1.3, 12.5, 5.8, RED)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "Interpretation & Analysis"; p.font.size = Pt(22); p.font.color.rgb = RED; p.font.bold = True
for t, c3 in [
    ("LIME has the LOWEST deletion AUC (0.054\u20130.152) \u2014 best at finding pixels that truly matter for the prediction", GREEN),
    ("Why? LIME removes entire superpixels (coherent regions), so each removal step has maximum impact on confidence.", LGRAY),
    ("", WHITE),
    ("Grad-CAM++ has higher deletion AUC (0.323 in-domain) \u2014 its coarse 7\u00d77 heatmap covers broadly, so individual pixel removal has less targeted effect", WHITE),
    ("", WHITE),
    ("Domain shift DESTROYS faithfulness: Real model drops from 0.323\u21920.125 on sketches (2.6\u00d7 worse)", YELLOW),
    ("The model's texture-based features are simply ABSENT in sketches \u2014 removing \"important\" sketch regions has less effect", LGRAY),
    ("", WHITE),
    ("Sketch model is ROBUST: deletion AUC varies only 0.205\u21920.194 across domains (5% change vs. 61% for real model)", GREEN),
    ("Shape-based features exist in both domains, so faithfulness is preserved during transfer", LGRAY),
]:
    p = tf.add_paragraph(); p.text = ("\u2022 " + t) if t else ""; p.font.size = Pt(15); p.font.color.rgb = c3; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Faithfulness — Insertion Table
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Faithfulness: Insertion AUC (\u2191 Higher is Better)")

tb(s, 0.5, 1.15, 12.3, 0.5, "Start from blank \u2192 reveal most important pixels first \u2192 confidence should RISE.\nHigher AUC = the explanation found the key features.", 16, LGRAY)

for i, (x, h) in enumerate(zip(dx, dh)):
    tb(s, x, 1.7, dw[i], 0.7, h, 15, ACCENT, True, PP_ALIGN.CENTER)

ins_data = [
    ("Grad-CAM++", ["0.537", "0.345", "0.453", "0.470"]),
    ("Grad-CAM",   ["0.548", "0.379", "0.473", "0.492"]),
    ("Integ. Grad.",["0.352", "0.206", "0.321", "0.272"]),
    ("LIME",       ["0.609", "0.501", "0.550", "0.571"]),
]
for j, (m, vals) in enumerate(ins_data):
    y = 2.5 + j * 0.55
    tb(s, dx[0], y, dw[0], 0.45, m, 16, WHITE, True)
    for k, v in enumerate(vals):
        fv = float(v)
        clr = GREEN if fv > 0.50 else (YELLOW if fv > 0.35 else (RED if fv < 0.25 else WHITE))
        tb(s, dx[k+1], y, dw[k+1], 0.45, v, 16, clr, fv > 0.55, PP_ALIGN.CENTER)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Insertion — Interpretation (separate slide)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Insertion: What Do These Numbers Mean?")

c2 = card(s, 0.4, 1.3, 12.5, 5.8, GREEN)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "Interpretation & Analysis"; p.font.size = Pt(22); p.font.color.rgb = GREEN; p.font.bold = True
for t, c3 in [
    ("LIME wins insertion across ALL conditions (0.501\u20130.609) \u2014 best at isolating the features that restore confidence", GREEN),
    ("Why? When top superpixels are revealed, the model sees COHERENT image patches (not scattered pixels) \u2192 faster confidence recovery", LGRAY),
    ("", WHITE),
    ("Grad-CAM/Grad-CAM++ are close (0.345\u20130.548). Their coarse heatmaps still identify generally important spatial regions.", WHITE),
    ("", WHITE),
    ("Integrated Gradients is worst (0.206\u20130.352). Pixel-level attributions are noisy \u2014 revealing scattered high-attribution pixels doesn't form meaningful image content.", WHITE),
    ("", WHITE),
    ("Cross-domain drop: Real model insertion falls from 0.609\u21920.501 (LIME) on sketches", YELLOW),
    ("The model recovers confidence SLOWER on unfamiliar visual styles \u2014 it can't extract useful features from the revealed sketch pixels", LGRAY),
]:
    p = tf.add_paragraph(); p.text = ("\u2022 " + t) if t else ""; p.font.size = Pt(15); p.font.color.rgb = c3; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Cross-Domain Consistency Table
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Cross-Domain Consistency (Cosine Similarity \u2014 \u2191 Higher is Better)")

tb(s, 0.5, 1.15, 12.3, 0.5, "Compare WHERE the model looks on real vs sketch images for the same class.\nRange: -1 to 1. Score of 1.0 = model focuses on identical regions across domains.", 16, LGRAY)

ch = ["Method", "Real Model\n(avg \u00b1 std)", "Sketch Model\n(avg \u00b1 std)", "Interpretation"]
ccx = [0.8, 4.0, 6.5, 9.0]
ccw = [3.0, 2.5, 2.5, 3.8]
for i, (x, h) in enumerate(zip(ccx, ch)):
    tb(s, x, 1.7, ccw[i], 0.7, h, 15, ACCENT, True, PP_ALIGN.CENTER)

cons_data = [
    ("Grad-CAM++", "0.963 \u00b1 0.028", "0.961 \u00b1 0.033", "Near-identical focus across domains", GREEN),
    ("Grad-CAM",   "0.953 \u00b1 0.035", "0.949 \u00b1 0.039", "Very consistent spatial attention", GREEN),
    ("LIME",       "0.871 \u00b1 0.042", "0.880 \u00b1 0.041", "Good consistency, some variation", YELLOW),
    ("Integ. Grad.","0.842 \u00b1 0.033", "0.895 \u00b1 0.024", "Moderate \u2014 pixel noise reduces agreement", YELLOW),
]
for j, (m, rv, sv, interp, clr) in enumerate(cons_data):
    y = 2.5 + j * 0.55
    tb(s, ccx[0], y, ccw[0], 0.45, m, 16, WHITE, True)
    tb(s, ccx[1], y, ccw[1], 0.45, rv, 16, clr, False, PP_ALIGN.CENTER)
    tb(s, ccx[2], y, ccw[2], 0.45, sv, 16, clr, False, PP_ALIGN.CENTER)
    tb(s, ccx[3], y, ccw[3], 0.45, interp, 14, LGRAY, False, PP_ALIGN.CENTER)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Consistency — Interpretation (separate slide)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Consistency: What Do These Numbers Mean?")

c2 = card(s, 0.4, 1.3, 12.5, 5.8, ACCENT)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]; p.text = "Interpretation & Analysis"; p.font.size = Pt(22); p.font.color.rgb = ACCENT; p.font.bold = True
for t, c3 in [
    ("All methods achieve > 0.84 \u2014 models look at broadly similar regions across real & sketch domains", WHITE),
    ("", WHITE),
    ("Grad-CAM++ is most consistent (0.96) \u2014 its coarse 7\u00d77 grid naturally aligns spatial focus. Minor pixel differences are averaged out.", WHITE),
    ("", WHITE),
    ("Despite the real model's 37-point accuracy DROP on sketches, it still LOOKS at the right regions (0.95 cosine)", YELLOW),
    ("It knows WHERE to look but can't extract useful features from what it sees \u2014 a previously unidentified failure mode", YELLOW),
    ("", WHITE),
    ("Sketch model IG consistency (0.895) > Real model IG (0.842) \u2014 shape-based learning produces more uniform attention patterns", WHITE),
    ("", WHITE),
    ("Key insight: spatial attention consistency \u2260 classification accuracy. A model can look at the right place but fail to interpret what it sees.", GREEN),
]:
    p = tf.add_paragraph(); p.text = ("\u2022 " + t) if t else ""; p.font.size = Pt(15); p.font.color.rgb = c3; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Representation Analysis — t-SNE by Domain
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Representation Analysis: t-SNE (Colored by Domain)")

tb(s, 0.5, 1.15, 12.3, 0.35, "t-SNE projects 2048-D feature vectors into 2D. Each dot = one test image. Color = domain.", 16, LGRAY)

tb(s, 0.8, 1.6, 5.8, 0.35, "Real Model", 22, RED, True, PP_ALIGN.CENTER)
img(s, f"{DOCS}/representations/real_tsne_domain.png", 0.8, 2.0, w=5.8)

tb(s, 6.8, 1.6, 5.8, 0.35, "Sketch Model", 22, GREEN, True, PP_ALIGN.CENTER)
img(s, f"{DOCS}/representations/sketch_tsne_domain.png", 6.8, 2.0, w=5.8)

c2 = card(s, 0.4, 5.8, 5.9, 1.4, RED)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06)
p = tf.paragraphs[0]; p.text = "Real Model: Domains SEPARATED"; p.font.size = Pt(16); p.font.color.rgb = RED; p.font.bold = True
p = tf.add_paragraph(); p.text = "Real & sketch form distinct clusters \u2192 the model\nlearned DIFFERENT features for each domain \u2192 fails\non sketches because it can't find texture features."; p.font.size = Pt(13); p.font.color.rgb = WHITE

c2 = card(s, 6.6, 5.8, 6.3, 1.4, GREEN)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06)
p = tf.paragraphs[0]; p.text = "Sketch Model: Domains MIXED"; p.font.size = Pt(16); p.font.color.rgb = GREEN; p.font.bold = True
p = tf.add_paragraph(); p.text = "Real & sketch images are interleaved \u2192 the model\nlearned SHARED shape-based features \u2192 transfers\nseamlessly because shape exists in both domains."; p.font.size = Pt(13); p.font.color.rgb = WHITE
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Representation Analysis — UMAP
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Representation Analysis: UMAP (Colored by Domain)")

tb(s, 0.5, 1.15, 12.3, 0.35, "UMAP preserves both local & global structure. Same pattern confirms: real model separates, sketch model mixes.", 16, LGRAY)

tb(s, 0.8, 1.6, 5.8, 0.35, "Real Model", 22, RED, True, PP_ALIGN.CENTER)
img(s, f"{DOCS}/representations/real_umap_domain.png", 0.8, 2.0, w=5.8)

tb(s, 6.8, 1.6, 5.8, 0.35, "Sketch Model", 22, GREEN, True, PP_ALIGN.CENTER)
img(s, f"{DOCS}/representations/sketch_umap_domain.png", 6.8, 2.0, w=5.8)

c2 = card(s, 0.4, 5.8, 12.5, 1.4, PURPLE)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.06)
p = tf.paragraphs[0]; p.text = "Both t-SNE and UMAP tell the same story"; p.font.size = Pt(18); p.font.color.rgb = PURPLE; p.font.bold = True
for t in [
    "Real model: domain-specific representation (texture-dependent) \u2192 two separate manifolds \u2192 transfer FAILS",
    "Sketch model: domain-agnostic representation (shape-based) \u2192 single unified manifold \u2192 transfer SUCCEEDS",
    "This is the mechanistic explanation for the asymmetric transfer: it's about HOW the model organizes features internally",
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = WHITE; p.space_after = Pt(2)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# t-SNE by Class
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Representation Analysis: t-SNE (Colored by Class)")

tb(s, 0.5, 1.15, 12.3, 0.35, "Same t-SNE plots, now colored by class label. Shows how well classes separate in feature space.", 16, LGRAY)

tb(s, 0.8, 1.6, 5.8, 0.35, "Real Model", 22, RED, True, PP_ALIGN.CENTER)
img(s, f"{DOCS}/representations/real_tsne_class.png", 0.8, 2.0, w=5.8)

tb(s, 6.8, 1.6, 5.8, 0.35, "Sketch Model", 22, GREEN, True, PP_ALIGN.CENTER)
img(s, f"{DOCS}/representations/sketch_tsne_class.png", 6.8, 2.0, w=5.8)

c2 = card(s, 0.4, 5.8, 12.5, 1.4, ACCENT)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.06)
p = tf.paragraphs[0]; p.text = "Class-level Structure"; p.font.size = Pt(18); p.font.color.rgb = ACCENT; p.font.bold = True
for t in [
    "Both models form class clusters, but the sketch model's clusters contain BOTH domain types \u2014 confirming domain-agnostic class learning",
    "Real model's class clusters are split by domain \u2014 e.g., 'real tiger' and 'sketch tiger' land in different areas despite being the same class",
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = WHITE; p.space_after = Pt(2)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Method Comparison Summary
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "XAI Method Comparison: No Single Winner")

tb(s, 0.5, 1.15, 12.3, 0.3, "Best value in each column highlighted. All values from real model, in-domain (real\u2192real).", 15, LGRAY)

comp_h = ["Method", "Stability\n(SSIM \u2191)", "Deletion\n(AUC \u2193)", "Insertion\n(AUC \u2191)", "Consistency\n(Cosine \u2191)", "Speed", "Resolution"]
comp_x = [0.5, 2.8, 4.8, 6.8, 8.8, 10.8, 12.0]
comp_w = [2.2, 2.0, 2.0, 2.0, 2.0, 1.2, 1.2]
for i, (x, h) in enumerate(zip(comp_x, comp_h)):
    tb(s, x, 1.5, comp_w[i], 0.7, h, 14, ACCENT, True, PP_ALIGN.CENTER)

comp_rows = [
    ("Grad-CAM", [("0.942", WHITE), ("0.294", WHITE), ("0.548", WHITE), ("0.953", WHITE), ("Fast", LGRAY), ("7\u00d77", LGRAY)]),
    ("Grad-CAM++", [("0.946", GREEN), ("0.323", GREEN), ("0.537", WHITE), ("0.963", GREEN), ("Fast", LGRAY), ("7\u00d77", LGRAY)]),
    ("Integ. Grad.", [("0.563", RED), ("0.277", WHITE), ("0.352", RED), ("0.842", RED), ("Slow", LGRAY), ("224\u00b2", LGRAY)]),
    ("LIME", [("0.471", RED), ("0.152", YELLOW), ("0.609", GREEN), ("0.871", YELLOW), ("Slowest", LGRAY), ("Superpx", LGRAY)]),
]
for j, (m, vals) in enumerate(comp_rows):
    y = 2.3 + j * 0.5
    tb(s, comp_x[0], y, comp_w[0], 0.45, m, 15, WHITE, True)
    for k, (v, clr) in enumerate(vals):
        tb(s, comp_x[k+1], y, comp_w[k+1], 0.45, v, 15, clr, clr == GREEN, PP_ALIGN.CENTER)

c2 = card(s, 0.4, 4.5, 12.5, 2.7, PURPLE)
tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; p.text = "Recommendations by Use Case"; p.font.size = Pt(20); p.font.color.rgb = PURPLE; p.font.bold = True
for t, c3 in [
    ("Grad-CAM++: Best default choice \u2014 wins stability, deletion faithfulness, and consistency. Fast & reliable.", GREEN),
    ("Grad-CAM: Very close runner-up. Simpler to implement. Choose for quick, reliable spatial explanations.", WHITE),
    ("LIME: Best for causal analysis \u2014 wins insertion (0.609). Use when you need to identify which regions truly CAUSE the prediction.", ORANGE),
    ("Integrated Gradients: Use when mathematical rigor is required (satisfies sensitivity + completeness axioms).", RED),
    ("No single method wins all axes \u2014 choose based on your priority.", YELLOW),
]:
    p = tf.add_paragraph(); p.text = "\u2022 " + t; p.font.size = Pt(14); p.font.color.rgb = c3; p.space_after = Pt(3)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Key Takeaways
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Key Takeaways")
takeaways = [
    ("1", "Texture Bias in CNNs", "Real-trained model drops 37.5 points on sketches (texture-dependent). Sketch-trained model generalizes\nto both domains (shape-based). Training on abstract data may produce more robust models.", ORANGE),
    ("2", "Explanations \u2260 Accuracy", "Even with 83% cross-domain accuracy, explanation faithfulness degrades 2.6\u00d7 (deletion AUC: 0.323\u21920.125).\nA correct prediction with a wrong explanation is dangerous in safety-critical applications.", RED),
    ("3", "Stability\u2013Granularity Trade-off", "CAM methods: stable (SSIM > 0.92) but coarse 7\u00d77. Pixel methods (IG): detailed 224\u00d7224 but SSIM ~ 0.55.\nNo free lunch \u2014 choose your resolution vs. reliability trade-off.", YELLOW),
    ("4", "No Single Best XAI Method", "Grad-CAM++ wins 3 of 4 axes but LIME wins insertion faithfulness (0.609 vs 0.537).\nChoose your method based on your evaluation priority.", GREEN),
    ("5", "Spatial Attention Survives Domain Shift", "All methods maintain > 0.84 cosine consistency across domains. Models know WHERE to look\nbut may lack features to interpret what they see \u2014 a previously unidentified failure mode.", ACCENT),
]
for i, (num, ttl, desc, color) in enumerate(takeaways):
    top = 1.25 + i * 1.2
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(top + 0.05), Inches(0.45), Inches(0.45))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    tf = sh.text_frame; p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(18); p.font.color.rgb = DARK_BG; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tb(s, 1.3, top, 5.0, 0.4, ttl, 22, color, True)
    tb(s, 1.3, top + 0.4, 11.5, 0.8, desc, 15, WHITE)
snum(s, slide_counter[0]); slide_counter[0] += 1

# ================================================================
# Future Work & Questions
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s, "Future Work & Questions")
bullets(s, 0.8, 1.3, 5.8, 4.5, [
    ("Future Directions", 0, ORANGE, True),
    ("Extend to more DomainNet domains (clipart, painting, quickdraw)", 1, WHITE),
    ("Test with Vision Transformers (ViT) + attention-based explanations", 1, WHITE),
    ("Apply domain adaptation (DANN, MMD) and measure explanation preservation", 1, WHITE),
    ("Add concept-based methods (TCAV) for higher-level explanations", 1, WHITE),
    ("Conduct human evaluation studies for explanation quality", 1, WHITE),
    ("Apply to medical imaging, autonomous driving", 1, WHITE),
    ("", 0, None),
    ("Broader Impact", 0, ORANGE, True),
    ("Framework is method-agnostic \u2014 can evaluate any new XAI technique", 1, WHITE),
    ("Results inform when to trust/distrust model explanations", 1, WHITE),
    ("Practical guide: which XAI method for which scenario", 1, WHITE),
], 19)
c2 = card(s, 7.0, 1.5, 5.5, 4.5, ACCENT)
tf = c2.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = ""; p.font.size = Pt(10)
p = tf.add_paragraph(); p.text = "Thank You!"; p.font.size = Pt(42); p.font.color.rgb = ACCENT; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""; p.font.size = Pt(20)
p = tf.add_paragraph(); p.text = "Questions?"; p.font.size = Pt(34); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""; p.font.size = Pt(20)
p = tf.add_paragraph(); p.text = "Prerak Patel"; p.font.size = Pt(22); p.font.color.rgb = ORANGE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "Florida Institute of Technology"; p.font.size = Pt(16); p.font.color.rgb = LGRAY; p.alignment = PP_ALIGN.CENTER
snum(s, slide_counter[0]); slide_counter[0] += 1

# Save
out = f"{BASE}/docs/XAI_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
